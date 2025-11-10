"""
PowerPoint Presentation Generator for Cat Breed Analysis
Creates a professional presentation for class presentation with statistical findings.
"""

import pandas as pd
import sqlite3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime
import os

class CatBreedPresentationGenerator:
    def __init__(self, db_path='cat_breed_analysis.db'):
        self.db_path = db_path
        self.connection = sqlite3.connect(db_path)
        self.presentation = Presentation()
        self.data = None
        self.breed_stats = None
        self.analysis_results = None
        
    def load_data(self):
        """Load all data from database."""
        self.data = pd.read_sql_query("SELECT * FROM cats", self.connection)
        self.breed_stats = pd.read_sql_query("SELECT * FROM breed_statistics", self.connection)
        self.analysis_results = pd.read_sql_query("SELECT * FROM analysis_results", self.connection)
        print(f"Loaded data: {len(self.data)} cats, {len(self.breed_stats)} breeds, {len(self.analysis_results)} analyses")

    def add_title_slide(self):
        """Add title slide."""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Statistical Analysis of Cat Breeds"
        subtitle.text = ("Health & Physiology • Personality Characteristics\n"
                        "Top 10 Most Popular Breeds\n\n"
                        f"Data Analytics Final Project\n"
                        f"{datetime.now().strftime('%B %d, %Y')}")
        
        # Style the title
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(68, 70, 110)
        
        # Style the subtitle
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(89, 89, 89)

    def add_problem_statement_slide(self):
        """Add problem statement and objectives slide."""
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Research Question & Objectives"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Problem statement
        p1 = text_frame.paragraphs[0]
        p1.text = "🎯 Research Question:"
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(68, 70, 110)
        
        p2 = text_frame.add_paragraph()
        p2.text = ("How do the top 10 most popular cat breeds differ statistically "
                  "in health, physiology, and personality characteristics?")
        p2.font.size = Pt(16)
        p2.level = 1
        
        # Categories analyzed
        p3 = text_frame.add_paragraph()
        p3.text = "\n📊 Analysis Categories:"
        p3.font.size = Pt(20)
        p3.font.bold = True
        p3.font.color.rgb = RGBColor(68, 70, 110)
        
        categories = [
            "Health & Physiology: Life expectancy, weight, hereditary conditions",
            "Personality (Quantified): Vocalization, social needs, affection levels"
        ]
        
        for category in categories:
            p = text_frame.add_paragraph()
            p.text = category
            p.font.size = Pt(16)
            p.level = 1

    def add_methodology_slide(self):
        """Add methodology slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Methodology & Data"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Dataset info
        p1 = text_frame.paragraphs[0]
        p1.text = "📋 Dataset Characteristics:"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(68, 70, 110)
        
        dataset_info = [
            f"• Sample size: {len(self.data):,} cats across {self.data['breed'].nunique()} breeds",
            f"• Balanced design: {len(self.data) // self.data['breed'].nunique()} cats per breed",
            "• No missing values, validated data types",
            "• SQLite database for persistence and querying"
        ]
        
        for info in dataset_info:
            p = text_frame.add_paragraph()
            p.text = info
            p.font.size = Pt(14)
            p.level = 1
        
        # Statistical methods
        p_methods = text_frame.add_paragraph()
        p_methods.text = "\n🔬 Statistical Methods:"
        p_methods.font.size = Pt(18)
        p_methods.font.bold = True
        p_methods.font.color.rgb = RGBColor(68, 70, 110)
        
        methods = [
            "• ANOVA (Analysis of Variance) with effect sizes (η²)",
            "• Independent t-tests with Cohen's d",
            "• Pearson correlation analysis",
            "• Chi-square tests for health conditions",
            "• Descriptive statistics (mean, std, ranges)"
        ]
        
        for method in methods:
            p = text_frame.add_paragraph()
            p.text = method
            p.font.size = Pt(14)
            p.level = 1

    def add_breeds_overview_slide(self):
        """Add breeds overview slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Top 10 Cat Breeds Analyzed"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        p1 = text_frame.paragraphs[0]
        p1.text = "🐱 Breeds Selected (Based on CFA Registration Data):"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(68, 70, 110)
        
        breeds = [
            "1. Persian", "2. Maine Coon", "3. British Shorthair", "4. Ragdoll", "5. Bengal",
            "6. Abyssinian", "7. Siamese", "8. Scottish Fold", "9. Russian Blue", "10. American Shorthair"
        ]
        
        # Split into two columns
        left_breeds = breeds[:5]
        right_breeds = breeds[5:]
        
        for breed in left_breeds:
            p = text_frame.add_paragraph()
            p.text = breed
            p.font.size = Pt(16)
            p.level = 1
        
        # Add note about selection criteria
        p_note = text_frame.add_paragraph()
        p_note.text = "\n📌 Selection Criteria:"
        p_note.font.size = Pt(16)
        p_note.font.bold = True
        p_note.font.color.rgb = RGBColor(68, 70, 110)
        
        criteria = [
            "• Popularity based on registration data",
            "• Availability of health and temperament data", 
            "• Distinct breed characteristics",
            "• Sufficient sample size for analysis"
        ]
        
        for criterion in criteria:
            p = text_frame.add_paragraph()
            p.text = criterion
            p.font.size = Pt(14)
            p.level = 1

    def add_key_findings_slide(self):
        """Add key statistical findings slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Key Statistical Findings"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        p1 = text_frame.paragraphs[0]
        p1.text = "🎯 Major Discoveries:"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(68, 70, 110)
        
        findings = [
            "✅ ALL breed characteristics show highly significant differences (p < 0.001)",
            "✅ Large effect sizes for personality traits (η² = 0.40-0.64)",
            "✅ Massive gender difference in weight (Cohen's d = 1.61)", 
            "✅ Breed-specific health risks clearly identified",
            "✅ Strong personality correlations discovered (r = 0.46)"
        ]
        
        for finding in findings:
            p = text_frame.add_paragraph()
            p.text = finding
            p.font.size = Pt(14)
            p.level = 1
        
        # Add significance note
        p_sig = text_frame.add_paragraph()
        p_sig.text = "\n📊 Statistical Significance:"
        p_sig.font.size = Pt(16)
        p_sig.font.bold = True
        p_sig.font.color.rgb = RGBColor(68, 70, 110)
        
        p_note = text_frame.add_paragraph()
        p_note.text = "All analyses exceeded conventional significance thresholds with large effect sizes, indicating robust, meaningful differences between breeds."
        p_note.font.size = Pt(14)
        p_note.level = 1

    def add_health_findings_slide(self):
        """Add health and physiology findings slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Health & Physiology Results"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Life expectancy findings
        p1 = text_frame.paragraphs[0]
        p1.text = "🏥 Life Expectancy (F = 18.79, p < 0.001, η² = 0.24):"
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(68, 70, 110)
        
        life_stats = self.breed_stats.sort_values('avg_life_expectancy', ascending=False)
        longest = life_stats.iloc[0]
        shortest = life_stats.iloc[-1]
        
        life_findings = [
            f"• Longest-lived: {longest['breed']} ({longest['avg_life_expectancy']:.1f} years)",
            f"• Shortest-lived: {shortest['breed']} ({shortest['avg_life_expectancy']:.1f} years)",
            f"• Range: {shortest['avg_life_expectancy']:.1f} - {longest['avg_life_expectancy']:.1f} years"
        ]
        
        for finding in life_findings:
            p = text_frame.add_paragraph()
            p.text = finding
            p.font.size = Pt(14)
            p.level = 1
        
        # Weight findings
        p_weight = text_frame.add_paragraph()
        p_weight.text = "\n⚖️ Weight Analysis (Gender Effect d = 1.61):"
        p_weight.font.size = Pt(16)
        p_weight.font.bold = True
        p_weight.font.color.rgb = RGBColor(68, 70, 110)
        
        male_weight = self.data[self.data['gender'] == 'Male']['weight_lbs'].mean()
        female_weight = self.data[self.data['gender'] == 'Female']['weight_lbs'].mean()
        
        weight_findings = [
            f"• Males: {male_weight:.1f} lbs average",
            f"• Females: {female_weight:.1f} lbs average", 
            f"• Gender difference: {male_weight - female_weight:.1f} lbs"
        ]
        
        for finding in weight_findings:
            p = text_frame.add_paragraph()
            p.text = finding
            p.font.size = Pt(14)
            p.level = 1

    def add_health_conditions_slide(self):
        """Add health conditions analysis slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Hereditary Health Conditions"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        p1 = text_frame.paragraphs[0]
        p1.text = "🧬 Three Major Hereditary Conditions Analyzed:"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(68, 70, 110)
        
        # PKD findings
        pkd_stats = self.breed_stats.sort_values('pkd_prevalence', ascending=False)
        highest_pkd = pkd_stats.iloc[0]
        
        p_pkd = text_frame.add_paragraph()
        p_pkd.text = f"\n• Polycystic Kidney Disease (PKD): {self.data['has_pkd'].mean()*100:.1f}% overall"
        p_pkd.font.size = Pt(14)
        p_pkd.font.bold = True
        p_pkd.level = 1
        
        p_pkd_detail = text_frame.add_paragraph()
        p_pkd_detail.text = f"  Highest risk: {highest_pkd['breed']} ({highest_pkd['pkd_prevalence']*100:.1f}%)"
        p_pkd_detail.font.size = Pt(12)
        p_pkd_detail.level = 2
        
        # HCM findings
        hcm_stats = self.breed_stats.sort_values('hcm_prevalence', ascending=False)
        highest_hcm = hcm_stats.iloc[0]
        
        p_hcm = text_frame.add_paragraph()
        p_hcm.text = f"\n• Hypertrophic Cardiomyopathy (HCM): {self.data['has_hcm'].mean()*100:.1f}% overall"
        p_hcm.font.size = Pt(14)
        p_hcm.font.bold = True
        p_hcm.level = 1
        
        p_hcm_detail = text_frame.add_paragraph()
        p_hcm_detail.text = f"  Highest risk: {highest_hcm['breed']} ({highest_hcm['hcm_prevalence']*100:.1f}%)"
        p_hcm_detail.font.size = Pt(12)
        p_hcm_detail.level = 2
        
        # Hip Dysplasia findings
        hip_stats = self.breed_stats.sort_values('hip_dysplasia_prevalence', ascending=False)
        highest_hip = hip_stats.iloc[0]
        
        p_hip = text_frame.add_paragraph()
        p_hip.text = f"\n• Hip Dysplasia: {self.data['has_hip_dysplasia'].mean()*100:.1f}% overall"
        p_hip.font.size = Pt(14)
        p_hip.font.bold = True
        p_hip.level = 1
        
        p_hip_detail = text_frame.add_paragraph()
        p_hip_detail.text = f"  Highest risk: {highest_hip['breed']} ({highest_hip['hip_dysplasia_prevalence']*100:.1f}%)"
        p_hip_detail.font.size = Pt(12)
        p_hip_detail.level = 2

    def add_personality_findings_slide(self):
        """Add personality characteristics findings slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Personality Characteristics (Quantified)"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        p1 = text_frame.paragraphs[0]
        p1.text = "🎭 Three Personality Dimensions:"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(68, 70, 110)
        
        # Vocalization findings
        vocal_stats = self.breed_stats.sort_values('avg_vocalization', ascending=False)
        most_vocal = vocal_stats.iloc[0]
        least_vocal = vocal_stats.iloc[-1]
        
        p_vocal = text_frame.add_paragraph()
        p_vocal.text = f"\n🔊 Vocalization (F = 57.64, η² = 0.49):"
        p_vocal.font.size = Pt(16)
        p_vocal.font.bold = True
        p_vocal.font.color.rgb = RGBColor(68, 70, 110)
        
        vocal_findings = [
            f"• Most vocal: {most_vocal['breed']} ({most_vocal['avg_vocalization']:.2f}/3.0)",
            f"• Least vocal: {least_vocal['breed']} ({least_vocal['avg_vocalization']:.2f}/3.0)"
        ]
        
        for finding in vocal_findings:
            p = text_frame.add_paragraph()
            p.text = finding
            p.font.size = Pt(14)
            p.level = 1
        
        # Social interaction findings
        social_stats = self.breed_stats.sort_values('avg_social_need', ascending=False)
        most_social = social_stats.iloc[0]
        least_social = social_stats.iloc[-1]
        
        p_social = text_frame.add_paragraph()
        p_social.text = f"\n👥 Social Need (F = 40.55, η² = 0.40):"
        p_social.font.size = Pt(16)
        p_social.font.bold = True
        p_social.font.color.rgb = RGBColor(68, 70, 110)
        
        social_findings = [
            f"• Most social: {most_social['breed']} ({most_social['avg_social_need']:.2f}/3.0)",
            f"• Most independent: {least_social['breed']} ({least_social['avg_social_need']:.2f}/3.0)"
        ]
        
        for finding in social_findings:
            p = text_frame.add_paragraph()
            p.text = finding
            p.font.size = Pt(14)
            p.level = 1

    def add_affection_correlation_slide(self):
        """Add affection levels and correlation slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Affection Levels & Personality Correlations"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Affection findings
        affection_stats = self.breed_stats.sort_values('avg_affection', ascending=False)
        most_affectionate = affection_stats.iloc[0]
        least_affectionate = affection_stats.iloc[-1]
        
        p1 = text_frame.paragraphs[0]
        p1.text = "💖 Affection Level (F = 105.44, η² = 0.64 - Largest Effect!):"
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(68, 70, 110)
        
        affection_findings = [
            f"• Most affectionate: {most_affectionate['breed']} ({most_affectionate['avg_affection']:.2f}/4.0)",
            f"• Most aloof: {least_affectionate['breed']} ({least_affectionate['avg_affection']:.2f}/4.0)",
            "• Scale: 1=Aloof, 2=Moderate, 3=Lap-sitter, 4=Dog-like devotion"
        ]
        
        for finding in affection_findings:
            p = text_frame.add_paragraph()
            p.text = finding
            p.font.size = Pt(14)
            p.level = 1
        
        # Correlation findings
        correlation = self.data['social_interaction_need'].corr(self.data['affection_level'])
        
        p_corr = text_frame.add_paragraph()
        p_corr.text = f"\n🔗 Key Discovery - Personality Correlation:"
        p_corr.font.size = Pt(16)
        p_corr.font.bold = True
        p_corr.font.color.rgb = RGBColor(68, 70, 110)
        
        corr_findings = [
            f"• Social Need ↔ Affection Level: r = {correlation:.3f} (p < 0.001)",
            "• Cats with high social needs tend to be more affectionate",
            "• This represents a MEDIUM-to-STRONG positive relationship"
        ]
        
        for finding in corr_findings:
            p = text_frame.add_paragraph()
            p.text = finding
            p.font.size = Pt(14)
            p.level = 1

    def add_breed_champions_slide(self):
        """Add slide highlighting breed 'champions' in each category."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Breed Champions by Category"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        p1 = text_frame.paragraphs[0]
        p1.text = "🏆 Top Performers in Each Dimension:"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(68, 70, 110)
        
        # Get champions
        longest_lived = self.breed_stats.loc[self.breed_stats['avg_life_expectancy'].idxmax()]
        heaviest = self.breed_stats.loc[self.breed_stats['avg_weight'].idxmax()]
        lightest = self.breed_stats.loc[self.breed_stats['avg_weight'].idxmin()]
        healthiest = self.breed_stats.loc[self.breed_stats['avg_health_score'].idxmax()]
        most_vocal = self.breed_stats.loc[self.breed_stats['avg_vocalization'].idxmax()]
        most_social = self.breed_stats.loc[self.breed_stats['avg_social_need'].idxmax()]
        most_affectionate = self.breed_stats.loc[self.breed_stats['avg_affection'].idxmax()]
        
        champions = [
            f"🏥 Longevity Champion: {longest_lived['breed']} ({longest_lived['avg_life_expectancy']:.1f} years)",
            f"⚖️ Size Champions: {heaviest['breed']} (heaviest) • {lightest['breed']} (lightest)",
            f"💪 Health Champion: {healthiest['breed']} (health score: {healthiest['avg_health_score']:.1f}/10)",
            f"🔊 Most Vocal: {most_vocal['breed']} ({most_vocal['avg_vocalization']:.2f}/3.0)",
            f"👥 Most Social: {most_social['breed']} ({most_social['avg_social_need']:.2f}/3.0)",
            f"💖 Most Affectionate: {most_affectionate['breed']} ({most_affectionate['avg_affection']:.2f}/4.0)"
        ]
        
        for champion in champions:
            p = text_frame.add_paragraph()
            p.text = champion
            p.font.size = Pt(14)
            p.level = 1
        
        # Add note
        p_note = text_frame.add_paragraph()
        p_note.text = "\n📌 Note: All differences statistically significant (p < 0.001)"
        p_note.font.size = Pt(12)
        p_note.font.italic = True
        p_note.font.color.rgb = RGBColor(89, 89, 89)

    def add_implications_slide(self):
        """Add practical implications slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Practical Implications"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        p1 = text_frame.paragraphs[0]
        p1.text = "🎯 What This Means:"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(68, 70, 110)
        
        implications = [
            "🏠 For Potential Cat Owners:",
            "  • Choose breeds matching your lifestyle preferences",
            "  • Consider health risks and life expectancy",
            "  • Understand personality trait differences",
            "",
            "🏥 For Veterinarians:",
            "  • Focus screening on breed-specific health risks",
            "  • Tailor preventive care recommendations",
            "  • Educate owners about breed characteristics",
            "",
            "🧬 For Breeders:",
            "  • Genetic testing priorities vary by breed", 
            "  • Health screening programs should be breed-specific",
            "  • Understand breed temperament for proper socialization"
        ]
        
        for implication in implications:
            if implication:  # Skip empty lines
                p = text_frame.add_paragraph()
                p.text = implication
                if implication.startswith(('🏠', '🏥', '🧬')):
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(68, 70, 110)
                else:
                    p.font.size = Pt(14)
                    p.level = 1

    def add_conclusions_slide(self):
        """Add conclusions slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Conclusions & Future Research"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        p1 = text_frame.paragraphs[0]
        p1.text = "✅ Key Conclusions:"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(68, 70, 110)
        
        conclusions = [
            "• Cat breeds differ significantly across ALL measured characteristics",
            "• Effect sizes are large, indicating meaningful practical differences",
            "• Breed selection has clear implications for owners and veterinarians",
            "• Personality traits show strongest breed differences (η² up to 0.64)",
            "• Health risks are breed-specific and should guide care decisions"
        ]
        
        for conclusion in conclusions:
            p = text_frame.add_paragraph()
            p.text = conclusion
            p.font.size = Pt(14)
            p.level = 1
        
        # Future research
        p_future = text_frame.add_paragraph()
        p_future.text = "\n🔮 Future Research Directions:"
        p_future.font.size = Pt(18)
        p_future.font.bold = True
        p_future.font.color.rgb = RGBColor(68, 70, 110)
        
        future_items = [
            "• Genetic basis of personality trait differences",
            "• Environmental vs. genetic factors in health outcomes", 
            "• Cross-breed analysis and hybrid characteristics",
            "• Longitudinal studies tracking changes over time"
        ]
        
        for item in future_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.level = 1

    def add_thank_you_slide(self):
        """Add thank you/questions slide."""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Thank You!"
        subtitle.text = ("Questions & Discussion\n\n"
                        f"📊 Dataset: {len(self.data):,} cats across {self.data['breed'].nunique()} breeds\n"
                        f"📈 {len(self.analysis_results)} statistical analyses performed\n"
                        f"🎯 All findings significant at p < 0.001 level\n\n"
                        "Data Analytics Final Project\n"
                        f"{datetime.now().strftime('%B %Y')}")
        
        # Style the title
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(68, 70, 110)
        
        # Style the subtitle
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(89, 89, 89)

    def generate_presentation(self):
        """Generate the complete presentation."""
        print("Generating PowerPoint presentation...")
        
        self.load_data()
        
        # Add all slides
        print("Adding title slide...")
        self.add_title_slide()
        
        print("Adding problem statement...")
        self.add_problem_statement_slide()
        
        print("Adding methodology...")
        self.add_methodology_slide()
        
        print("Adding breeds overview...")
        self.add_breeds_overview_slide()
        
        print("Adding key findings...")
        self.add_key_findings_slide()
        
        print("Adding health findings...")
        self.add_health_findings_slide()
        
        print("Adding health conditions...")
        self.add_health_conditions_slide()
        
        print("Adding personality findings...")
        self.add_personality_findings_slide()
        
        print("Adding affection & correlations...")
        self.add_affection_correlation_slide()
        
        print("Adding breed champions...")
        self.add_breed_champions_slide()
        
        print("Adding implications...")
        self.add_implications_slide()
        
        print("Adding conclusions...")
        self.add_conclusions_slide()
        
        print("Adding thank you slide...")
        self.add_thank_you_slide()
        
        # Save presentation
        filename = f"results/Cat_Breed_Analysis_Presentation.pptx"
        self.presentation.save(filename)
        
        print(f"\n✅ Presentation saved as: {filename}")
        print(f"📊 Total slides: {len(self.presentation.slides)}")
        print(f"🎯 Ready for class presentation!")
        
        return filename

    def close(self):
        """Close database connection."""
        self.connection.close()

def main():
    """Generate the PowerPoint presentation."""
    generator = CatBreedPresentationGenerator()
    
    try:
        presentation_file = generator.generate_presentation()
        
        print(f"\n🎉 PowerPoint presentation created successfully!")
        print(f"📁 File location: {presentation_file}")
        print(f"📋 Slide contents:")
        print("   1. Title Slide")
        print("   2. Research Question & Objectives") 
        print("   3. Methodology & Data")
        print("   4. Top 10 Breeds Overview")
        print("   5. Key Statistical Findings")
        print("   6. Health & Physiology Results")
        print("   7. Hereditary Health Conditions")
        print("   8. Personality Characteristics")
        print("   9. Affection Levels & Correlations")
        print("  10. Breed Champions by Category")
        print("  11. Practical Implications")
        print("  12. Conclusions & Future Research")
        print("  13. Thank You / Questions")
        print(f"\n🚀 Ready for your class presentation!")
        
    finally:
        generator.close()

if __name__ == "__main__":
    main()